#!/usr/bin/env python
"""Generate a clean 15-slide PowerPoint for the Master Colloquium B presentation.

Run:  python scripts/build_pptx.py
Output: docs/presentation/presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ─── Design system ────────────────────────────────────────────────────────────
INK      = RGBColor(0x0F, 0x17, 0x2A)   # near-black slate (primary text)
ACCENT   = RGBColor(0x25, 0x63, 0xEB)   # blue-600 (accent)
ACCENTDK = RGBColor(0x1E, 0x29, 0x4B)   # deep navy (section dividers)
MUTED    = RGBColor(0x64, 0x74, 0x8B)   # slate-500 (secondary text)
LIGHT    = RGBColor(0xF1, 0xF5, 0xF9)   # slate-100 (fills)
LINE     = RGBColor(0xE2, 0xE8, 0xF0)   # slate-200 (borders)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x05, 0x96, 0x69)   # emerald (positive)

FONT   = "Segoe UI"
FONT_L = "Segoe UI Light"
FONT_SB = "Segoe UI Semibold"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.9)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

PAGE = {"n": 0}


def _set_font(run, size, color, name=FONT, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = name
    run.font.bold = bold
    run.font.italic = italic


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size, color, name=FONT, bold=False, italic=False,
         space_after=8, space_before=0, align=PP_ALIGN.LEFT, level=0, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    _set_font(r, size, color, name, bold, italic)
    return p


def bullet(tf, segments, size=18, space_after=10, level=0, first=False, bullet_color=ACCENT):
    """segments: list of (text, bold) run tuples on one line."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    # bullet marker
    m = p.add_run()
    m.text = ("•  " if level == 0 else "–  ")
    _set_font(m, size, bullet_color if level == 0 else MUTED, FONT, bold=True)
    for text, bold in segments:
        r = p.add_run()
        r.text = text
        _set_font(r, size, INK if level == 0 else MUTED, FONT, bold=bold)
    return p


def footer(slide, dark=False):
    PAGE["n"] += 1
    col = RGBColor(0xB6, 0xC2, 0xD9) if dark else MUTED
    tb, tf = textbox(slide, MARGIN, Inches(7.02), Inches(9), Inches(0.35))
    para(tf, "IaC vs Manual AWS Provisioning  ·  Master Colloquium B",
         10, col, first=True)
    tb2, tf2 = textbox(slide, Inches(11.6), Inches(7.02), Inches(1.6), Inches(0.35))
    para(tf2, str(PAGE["n"]), 10, col, align=PP_ALIGN.RIGHT, first=True)


def title_block(slide, kicker, title):
    # small accent bar
    rect(slide, MARGIN, Inches(0.62), Inches(0.55), Inches(0.09), fill=ACCENT)
    tb, tf = textbox(slide, MARGIN, Inches(0.78), Inches(11.5), Inches(0.4))
    para(tf, kicker.upper(), 13, ACCENT, name=FONT_SB, bold=True, first=True)
    tb2, tf2 = textbox(slide, MARGIN, Inches(1.15), Inches(11.5), Inches(0.9))
    para(tf2, title, 30, INK, name=FONT_SB, bold=True, first=True)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def body_area():
    return MARGIN, Inches(2.25), Inches(11.5), Inches(4.5)


# ─── Slide 1 — Title ──────────────────────────────────────────────────────────
def slide_title():
    s = add_slide()
    rect(s, 0, 0, EMU_W, EMU_H, fill=WHITE)
    # left accent band
    rect(s, 0, 0, Inches(0.35), EMU_H, fill=ACCENT)
    rect(s, MARGIN, Inches(1.9), Inches(0.8), Inches(0.12), fill=ACCENT)
    tb, tf = textbox(s, MARGIN, Inches(2.15), Inches(11.6), Inches(2.2))
    para(tf, "Infrastructure as Code", 44, INK, name=FONT_L, first=True, space_after=2)
    para(tf, "vs. Manual Provisioning", 44, INK, name=FONT_SB, bold=True, space_after=10)
    para(tf, "A comparative study of Terraform and AWS Console deployment",
         20, MUTED, space_after=0)
    tb2, tf2 = textbox(s, MARGIN, Inches(5.2), Inches(11.6), Inches(1.6))
    para(tf2, "Rafael Aza", 18, INK, name=FONT_SB, bold=True, first=True, space_after=2)
    para(tf2, "Master Colloquium B", 15, MUTED, space_after=2)
    para(tf2, "Supervisor: Prof. Dr.-Ing. Marcus Purat  ·  July 15, 2026", 15, MUTED)
    notes(s, "Greet the room. One-line hook: comparing building cloud infra by hand "
              "vs. as code. Mention it is a mid-project talk; measured results follow in the report. (~30s)")


# ─── Agenda ───────────────────────────────────────────────────────────────────
def slide_agenda():
    s = add_slide()
    title_block(s, "Overview", "Agenda")
    x, y, w, h = body_area()
    items = [
        ("1", "Context", "What is IaC, and why it matters"),
        ("2", "Target", "The problem & research question"),
        ("3", "Requirements", "Evaluation criteria & definitions"),
        ("4", "Functional Design", "The 3-tier architecture & tooling"),
        ("5", "Methodology & preliminary results", "Work in progress"),
        ("6", "Outlook", "Next steps & future work"),
    ]
    yy = y
    for num, head, sub in items:
        rect(s, x, yy + Inches(0.02), Inches(0.5), Inches(0.5), fill=LIGHT, shape=MSO_SHAPE.OVAL)
        tbn, tfn = textbox(s, x, yy + Inches(0.02), Inches(0.5), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
        para(tfn, num, 16, ACCENT, name=FONT_SB, bold=True, align=PP_ALIGN.CENTER, first=True)
        tb, tf = textbox(s, x + Inches(0.75), yy, Inches(10), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = head + "   "; _set_font(r, 18, INK, FONT_SB, bold=True)
        r2 = p.add_run(); r2.text = "— " + sub; _set_font(r2, 15, MUTED, FONT)
        yy += Inches(0.72)
    footer(s)
    notes(s, "Walk the six points quickly. Set expectation: focus on context, design, "
              "outlook today; full results in the report after feedback. (~30s)")


# ─── Section divider ──────────────────────────────────────────────────────────
def slide_section(num, title, subtitle=""):
    s = add_slide()
    rect(s, 0, 0, EMU_W, EMU_H, fill=ACCENTDK)
    rect(s, MARGIN, Inches(2.7), Inches(0.9), Inches(0.14), fill=ACCENT)
    tb, tf = textbox(s, MARGIN, Inches(2.95), Inches(11.5), Inches(2))
    para(tf, f"{num}", 20, RGBColor(0x8F,0xA8,0xE0), name=FONT_SB, bold=True, first=True, space_after=4)
    para(tf, title, 40, WHITE, name=FONT_SB, bold=True, space_after=6)
    if subtitle:
        para(tf, subtitle, 18, RGBColor(0xB6,0xC2,0xD9))
    footer(s, dark=True)
    return s


# ─── Context 1 ────────────────────────────────────────────────────────────────
def slide_context_what():
    s = add_slide()
    title_block(s, "1 · Context", "Two ways to build cloud infrastructure")
    x, y, w, h = body_area()
    colw = Inches(5.5)
    gap = Inches(0.5)
    # Card 1 — Click-Ops
    rect(s, x, y, colw, Inches(3.6), fill=LIGHT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb, tf = textbox(s, x + Inches(0.35), y + Inches(0.3), colw - Inches(0.7), Inches(3))
    para(tf, "Manual  ·  \u201cClick-Ops\u201d", 20, INK, name=FONT_SB, bold=True, first=True, space_after=10)
    bullet(tf, [("Create each resource by hand", False)], 16)
    bullet(tf, [("Click through the AWS Console", False)], 16)
    bullet(tf, [("Every VPC, server, firewall rule", False)], 16)
    bullet(tf, [("Manual, sequential, repetitive", False)], 16, space_after=0)
    # Card 2 — IaC
    x2 = x + colw + gap
    rect(s, x2, y, colw, Inches(3.6), fill=RGBColor(0xEA,0xF1,0xFE), line=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb2, tf2 = textbox(s, x2 + Inches(0.35), y + Inches(0.3), colw - Inches(0.7), Inches(3))
    para(tf2, "Infrastructure as Code", 20, ACCENT, name=FONT_SB, bold=True, first=True, space_after=10)
    bullet(tf2, [("Define infra as ", False), ("version-controlled code", True)], 16)
    bullet(tf2, [("Provision with ", False), ("one command", True)], 16)
    bullet(tf2, [("Declarative & state-aware", False)], 16)
    bullet(tf2, [("Tool: ", False), ("Terraform", True), (" (industry standard)", False)], 16, space_after=0)
    footer(s)
    notes(s, "Define Click-Ops and IaC clearly. Name Terraform as the de-facto standard. "
              "Keep plain and quick. (~1.5 min)")


# ─── Context 2 ────────────────────────────────────────────────────────────────
def slide_context_why():
    s = add_slide()
    title_block(s, "1 · Context", "Why it matters")
    x, y, w, h = body_area()
    tb, tf = textbox(s, x, y, Inches(7.2), Inches(4))
    para(tf, "The problems with doing it by hand:", 18, INK, name=FONT_SB, bold=True, first=True, space_after=12)
    bullet(tf, [("Doesn\u2019t scale", True), (" beyond a handful of resources", False)], 18)
    bullet(tf, [("Environments ", False), ("drift apart", True), (" over time", False)], 18)
    bullet(tf, [("No ", False), ("audit trail", True), (" of who changed what", False)], 18)
    bullet(tf, [("Error-prone", True), (" and hard to reproduce", False)], 18, space_after=16)
    para(tf, "Motivation: measure the difference with real data \u2014 not vendor claims.",
         17, ACCENT, name=FONT_SB, bold=True, italic=True)
    # resources side panel
    px = x + Inches(7.7)
    rect(s, px, y, Inches(3.8), Inches(3.7), fill=LIGHT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tbp, tfp = textbox(s, px + Inches(0.3), y + Inches(0.28), Inches(3.2), Inches(3.2))
    para(tfp, "RESOURCES USED", 13, MUTED, name=FONT_SB, bold=True, first=True, space_after=10)
    for item in ["AWS  ·  eu-central-1", "Terraform v1.15", "WSL2 Ubuntu, AWS CLI, Node.js", "AWS promotional credits", "Git / GitHub"]:
        para(tfp, item, 15, INK, space_after=9)
    footer(s)
    notes(s, "Three pains: scale, drift, no audit trail. Motivation is to MEASURE, not "
              "repeat marketing. Quickly list resources from the panel. (~1 min)")


# ─── Problem / Question ───────────────────────────────────────────────────────
def slide_problem():
    s = add_slide()
    title_block(s, "2 · Target", "The problem & research question")
    x, y, w, h = body_area()
    tb, tf = textbox(s, x, y, Inches(11.5), Inches(1.6))
    para(tf, "Provisioning a realistic, multi-tier environment by hand is:", 18, INK, first=True, space_after=10)
    p = tf.add_paragraph(); p.space_after = Pt(0)
    for i, (word) in enumerate(["Slow", "Inconsistent", "Hard to reproduce", "Opaque"]):
        r = p.add_run(); r.text = ("      " if i else "") + "●  "; _set_font(r, 16, ACCENT, FONT, bold=True)
        r2 = p.add_run(); r2.text = word; _set_font(r2, 18, INK, FONT_SB, bold=True)
    # research question card
    qy = y + Inches(1.85)
    rect(s, x, qy, Inches(11.5), Inches(1.9), fill=ACCENTDK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, qy, Inches(0.16), Inches(1.9), fill=ACCENT)
    tbq, tfq = textbox(s, x + Inches(0.5), qy + Inches(0.28), Inches(10.6), Inches(1.4), anchor=MSO_ANCHOR.MIDDLE)
    para(tfq, "RESEARCH QUESTION", 12, RGBColor(0x8F,0xA8,0xE0), name=FONT_SB, bold=True, first=True, space_after=8)
    para(tfq, "To what extent does IaC (Terraform) improve the deployment speed and "
              "operational consistency of AWS resources compared to manual provisioning?",
         21, WHITE, name=FONT_SB, bold=True)
    footer(s)
    notes(s, "State the problem in one breath. Then READ the research question slowly \u2014 "
              "it anchors the whole talk. Two measurable dimensions: speed & consistency. (~1.5 min)")


# ─── Requirements ─────────────────────────────────────────────────────────────
def slide_requirements():
    s = add_slide()
    title_block(s, "3 · Requirements", "Evaluation criteria & definitions")
    x, y, w, h = body_area()
    tb, tf = textbox(s, x, y - Inches(0.15), Inches(11.5), Inches(0.6))
    para(tf, "A production-shaped 3-tier system, deployed twice (manual & Terraform) under identical specs.",
         16, MUTED, italic=True, first=True)
    rows = [
        ("Criterion", "Definition", "How measured"),
        ("Deployment speed", "Time from start to app serving traffic", "Stopwatch  vs  time terraform apply"),
        ("Consistency / drift", "Reproducibility across runs", "Error count  vs  terraform plan = \u201cNo changes\u201d"),
        ("Auditability", "Ability to trace configuration", "Manual notes  vs  state file + Git"),
        ("Rollback", "Effort to tear the stack down", "Console deletion  vs  terraform destroy"),
    ]
    ty = y + Inches(0.55)
    tbl = s.shapes.add_table(len(rows), 3, x, ty, Inches(11.5), Inches(3.4)).table
    tbl.columns[0].width = Inches(2.9)
    tbl.columns[1].width = Inches(4.0)
    tbl.columns[2].width = Inches(4.6)
    for ci in range(3):
        _style_cell(tbl.cell(0, ci), rows[0][ci], header=True)
    for ri in range(1, len(rows)):
        for ci in range(3):
            _style_cell(tbl.cell(ri, ci), rows[ri][ci], header=False, first_col=(ci == 0))
    footer(s)
    notes(s, "Describe the target system in one line. Then walk the criteria table \u2014 "
              "this shows rigor; each criterion has a concrete measurement. Linger here. (~1.5 min)")


def _style_cell(cell, text, header=False, first_col=False):
    cell.margin_left = Inches(0.15)
    cell.margin_right = Inches(0.15)
    cell.margin_top = Inches(0.06)
    cell.margin_bottom = Inches(0.06)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENTDK if header else WHITE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    if header:
        _set_font(r, 14, WHITE, FONT_SB, bold=True)
    else:
        _set_font(r, 14, INK if first_col else MUTED, FONT, bold=first_col)


# ─── Architecture ─────────────────────────────────────────────────────────────
def slide_architecture():
    s = add_slide()
    title_block(s, "4 · Functional Design", "The 3-tier architecture")
    cx = Inches(5.95)         # center x for the stack (shifted right to fit left labels)
    boxw, boxh = Inches(5.4), Inches(0.92)
    x = cx - boxw / 2
    y0 = Inches(2.15)
    gap = Inches(0.34)
    tiers = [
        ("Internet  \u2192  Internet Gateway", MUTED, LIGHT, INK),
        ("Application Load Balancer", ACCENT, RGBColor(0xEA,0xF1,0xFE), ACCENT),
        ("Auto Scaling Group  ·  EC2 t3.micro  ·  Node.js", ACCENT, RGBColor(0xEA,0xF1,0xFE), ACCENT),
        ("RDS MySQL 8.0  (Multi-AZ)", ACCENT, RGBColor(0xEA,0xF1,0xFE), ACCENT),
    ]
    ys = []
    for i, (label, border, fill, txt) in enumerate(tiers):
        yy = y0 + i * (boxh + gap)
        ys.append(yy)
        rect(s, x, yy, boxw, boxh, fill=fill, line=border, line_w=Pt(1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb, tf = textbox(s, x, yy, boxw, boxh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, label, 16, txt, name=FONT_SB, bold=True, align=PP_ALIGN.CENTER, first=True)
    # arrows between boxes
    for i in range(len(tiers) - 1):
        ay = ys[i] + boxh
        conn = rect(s, cx - Inches(0.02), ay, Inches(0.04), gap, fill=ACCENT)
    # AZ annotation on the right
    tbz, tfz = textbox(s, x + boxw + Inches(0.35), y0, Inches(2.4), Inches(4))
    para(tfz, "Spans 2\nAvailability\nZones", 15, MUTED, name=FONT_SB, bold=True, first=True)
    # subnet annotation on the left
    labels = ["Public subnets", "Public subnets", "Private app subnets", "Private DB subnets"]
    for i, lb in enumerate(labels):
        tl, tfl = textbox(s, Inches(0.75), ys[i], Inches(2.25), boxh, anchor=MSO_ANCHOR.MIDDLE)
        para(tfl, lb, 12, MUTED, align=PP_ALIGN.RIGHT, first=True)
    footer(s)
    notes(s, "Talk top to bottom: Internet \u2192 Load Balancer \u2192 auto-scaling app \u2192 Multi-AZ DB. "
              "Stress two AZs for fault tolerance. \u2018Deliberately realistic, not a toy.\u2019 (~1.5 min)")


# ─── Design highlights ────────────────────────────────────────────────────────
def slide_design_highlights():
    s = add_slide()
    title_block(s, "4 · Functional Design", "Design highlights")
    x, y, w, h = body_area()
    tb, tf = textbox(s, x, y, Inches(11.5), Inches(2.4))
    bullet(tf, [("Network: ", True), ("custom VPC, 6 subnets across 2 AZs, Internet + NAT gateways", False)], 18, first=True)
    bullet(tf, [("Security: ", True), ("least-privilege chain \u2014 nothing reaches the database directly", False)], 18)
    bullet(tf, [("Compute: ", True), ("load-balanced, auto-scaling application tier", False)], 18)
    bullet(tf, [("Data: ", True), ("Multi-AZ RDS for fault tolerance", False)], 18)
    bullet(tf, [("Payload: ", True), ("a small app that queries the DB \u2014 proves the full chain works", False)], 18, space_after=18)
    # security chain strip
    cy = y + Inches(2.9)
    chain = ["Internet", "ALB :80", "EC2 :4000", "RDS :3306"]
    cw = Inches(2.3); cgap = Inches(0.55); cxx = x
    for i, node in enumerate(chain):
        fill = LIGHT if i == 0 else RGBColor(0xEA,0xF1,0xFE)
        bd = LINE if i == 0 else ACCENT
        rect(s, cxx, cy, cw, Inches(0.7), fill=fill, line=bd, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tbc, tfc = textbox(s, cxx, cy, cw, Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
        para(tfc, node, 15, INK, name=FONT_SB, bold=True, align=PP_ALIGN.CENTER, first=True)
        if i < len(chain) - 1:
            ta, tfa = textbox(s, cxx + cw, cy, cgap, Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
            para(tfa, "\u2192", 20, ACCENT, name=FONT_SB, bold=True, align=PP_ALIGN.CENTER, first=True)
        cxx += cw + cgap
    footer(s)
    notes(s, "Four tiers briefly, then the security chain: traffic only flows one way. "
              "Payload actually queries the DB, proving end-to-end. (~1 min)")


# ─── Terraform approach ───────────────────────────────────────────────────────
def slide_terraform():
    s = add_slide()
    title_block(s, "4 · Functional Design", "The Terraform approach")
    x, y, w, h = body_area()
    tb, tf = textbox(s, x, y, Inches(7.0), Inches(4))
    bullet(tf, [("Split into ", False), ("4 reusable modules", True), (":", False)], 18, first=True, space_after=6)
    para(tf, "     vpc  ·  security  ·  database  ·  compute", 16, MUTED, name=FONT_SB, bold=True, space_after=14)
    bullet(tf, [("One command", True), (" provisions everything in dependency order", False)], 18)
    bullet(tf, [("Every resource is code", True), (" \u2014 reviewable, repeatable, versioned", False)], 18, space_after=14)
    # testing callout
    px = x + Inches(7.5)
    rect(s, px, y, Inches(4.0), Inches(3.5), fill=RGBColor(0xEC,0xFD,0xF5), line=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tbp, tfp = textbox(s, px + Inches(0.32), y + Inches(0.3), Inches(3.4), Inches(3))
    para(tfp, "TESTED BEFORE DEPLOY", 13, GREEN, name=FONT_SB, bold=True, first=True, space_after=12)
    para(tfp, "\u2713  terraform validate", 16, INK, space_after=9)
    para(tfp, "\u2713  14 automated plan tests", 16, INK, space_after=9)
    para(tfp, "\u2713  application unit tests", 16, INK, space_after=9)
    para(tfp, "All green \u2014 no cloud needed", 14, MUTED, italic=True)
    footer(s)
    notes(s, "Four modules, one command, dependency order. Testing is the differentiator: "
              "validate + 14 plan tests + app unit tests, all before touching the cloud. (~1 min)")


# ─── Methodology ──────────────────────────────────────────────────────────────
def slide_methodology():
    s = add_slide()
    title_block(s, "5 · Methodology", "How the comparison is measured")
    x, y, w, h = body_area()
    steps = [
        ("Manual run", "Provision the full stack via the AWS Console. Stopwatch the process; log every error and rework."),
        ("Terraform run", "time terraform apply. Verify the app responds; re-run terraform plan to prove zero drift."),
        ("Compare", "Speed, consistency, auditability, rollback \u2014 side by side."),
    ]
    yy = y
    for i, (head, body) in enumerate(steps, 1):
        rect(s, x, yy, Inches(0.6), Inches(0.6), fill=ACCENT, shape=MSO_SHAPE.OVAL)
        tbn, tfn = textbox(s, x, yy, Inches(0.6), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
        para(tfn, str(i), 18, WHITE, name=FONT_SB, bold=True, align=PP_ALIGN.CENTER, first=True)
        tb, tf = textbox(s, x + Inches(0.9), yy - Inches(0.05), Inches(10.4), Inches(1.1))
        para(tf, head, 18, INK, name=FONT_SB, bold=True, first=True, space_after=3)
        para(tf, body, 15, MUTED)
        yy += Inches(1.25)
    tbf, tff = textbox(s, x, yy + Inches(0.05), Inches(11.4), Inches(0.5))
    para(tff, "Both environments are destroyed the same day \u2014 strict cost control.",
         15, ACCENT, name=FONT_SB, bold=True, italic=True, first=True)
    footer(s)
    notes(s, "Three steps: manual (stopwatch + errors), Terraform (timed apply + drift check), "
              "compare. Note same-day teardown for cost. (~1 min)")


# ─── Preliminary results ──────────────────────────────────────────────────────
def slide_results():
    s = add_slide()
    title_block(s, "5 · Preliminary Results", "Phase-1 pilot data")
    x, y, w, h = body_area()
    tb, tf = textbox(s, x, y - Inches(0.15), Inches(11.5), Inches(0.5))
    para(tf, "A simplified pilot (single VPC + EC2) validated the measurement method.",
         16, MUTED, italic=True, first=True)
    # big numbers
    ny = y + Inches(0.5)
    cards = [("Manual (Console)", "01:17:40", MUTED, LIGHT, LINE),
             ("Terraform", "00:20:00", ACCENT, RGBColor(0xEA,0xF1,0xFE), ACCENT)]
    cw = Inches(4.0); cgap = Inches(0.6); cxx = x
    for label, val, col, fill, bd in cards:
        rect(s, cxx, ny, cw, Inches(1.9), fill=fill, line=bd, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tbc, tfc = textbox(s, cxx, ny + Inches(0.28), cw, Inches(1.4), anchor=MSO_ANCHOR.TOP)
        para(tfc, label, 15, MUTED, name=FONT_SB, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=6)
        para(tfc, val, 40, col, name=FONT_SB, bold=True, align=PP_ALIGN.CENTER)
        cxx += cw + cgap
    # highlight
    hx = cxx
    rect(s, hx, ny, Inches(2.3), Inches(1.9), fill=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tbh, tfh = textbox(s, hx, ny + Inches(0.35), Inches(2.3), Inches(1.3), anchor=MSO_ANCHOR.MIDDLE)
    para(tfh, "~3.8\u00d7", 40, WHITE, name=FONT_SB, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
    para(tfh, "faster", 16, WHITE, align=PP_ALIGN.CENTER)
    tbn, tfn = textbox(s, x, ny + Inches(2.2), Inches(11.4), Inches(0.7))
    para(tfn, "Even at small scale, Terraform is decisively faster \u2014 the full 3-tier run is expected to widen the gap.",
         16, INK, italic=True, first=True)
    footer(s)
    notes(s, "Be honest: Phase-1 pilot to validate method. Headline ~3.8x faster (1h17 vs 20m). "
              "Expect the full 3-tier gap to be larger \u2014 manual complexity grows faster than code. (~1.5 min)")


# ─── Status ───────────────────────────────────────────────────────────────────
def slide_status():
    s = add_slide()
    title_block(s, "5 · Status", "Where the project stands")
    x, y, w, h = body_area()
    # done column
    rect(s, x, y, Inches(5.5), Inches(3.6), fill=RGBColor(0xEC,0xFD,0xF5), line=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb, tf = textbox(s, x + Inches(0.35), y + Inches(0.3), Inches(4.9), Inches(3))
    para(tf, "DONE", 14, GREEN, name=FONT_SB, bold=True, first=True, space_after=12)
    for t in ["Architecture designed", "Terraform code written (4 modules)", "Application built", "All tests passing", "On GitHub, version-controlled"]:
        para(tf, "\u2713  " + t, 16, INK, space_after=9)
    # next column
    x2 = x + Inches(6.0)
    rect(s, x2, y, Inches(5.5), Inches(3.6), fill=LIGHT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb2, tf2 = textbox(s, x2 + Inches(0.35), y + Inches(0.3), Inches(4.9), Inches(3))
    para(tf2, "NEXT", 14, ACCENT, name=FONT_SB, bold=True, first=True, space_after=12)
    for t in ["Run full 3-tier manual deployment", "Run full 3-tier Terraform deployment", "Capture real timings & drift data", "Comparative charts for the report"]:
        para(tf2, "\u2192  " + t, 16, INK, space_after=9)
    footer(s)
    notes(s, "Code and app are written, tested, version-controlled \u2014 ready to deploy. "
              "Remaining: run the full experiment and capture numbers for the report. (~30s)")


# ─── Outlook ──────────────────────────────────────────────────────────────────
def slide_outlook():
    s = slide_section("6 · Outlook", "Where this goes next")
    x = MARGIN
    tb, tf = textbox(s, x, Inches(4.6), Inches(11.5), Inches(2.2))
    items = [
        ("Complete the empirical run", "full 3-tier data for the report"),
        ("CI/CD integration", "Terraform on every git push"),
        ("Remote state & locking", "team-safe collaboration"),
        ("Policy as code", "automated security checks"),
        ("Cost analysis", "quantify operational savings"),
    ]
    p = tf.paragraphs[0]
    for i, (head, sub) in enumerate(items):
        pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pp.space_after = Pt(8)
        r = pp.add_run(); r.text = "\u2192  "; _set_font(r, 17, ACCENT, FONT, bold=True)
        r2 = pp.add_run(); r2.text = head + "  "; _set_font(r2, 17, WHITE, FONT_SB, bold=True)
        r3 = pp.add_run(); r3.text = "\u2014 " + sub; _set_font(r3, 16, RGBColor(0xB6,0xC2,0xD9), FONT)
    notes(s, "Say 2\u20133 aloud (don\u2019t read all): CI/CD, remote state, policy-as-code, cost. "
              "Frame as natural extensions. (~1 min)")


# ─── Summary ──────────────────────────────────────────────────────────────────
def slide_summary():
    s = add_slide()
    title_block(s, "Summary", "Key takeaways")
    x, y, w, h = body_area()
    tb, tf = textbox(s, x, y, Inches(11.5), Inches(4))
    points = [
        "IaC replaces error-prone manual clicking with version-controlled, repeatable code",
        "Designed & built a production-grade 3-tier AWS architecture, deployable both ways",
        "Defined clear, measurable criteria: speed, consistency, auditability, rollback",
        "Pilot data already shows a decisive speed advantage for Terraform",
        "Full empirical comparison follows in the report",
    ]
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run(); r.text = "\u2014  "; _set_font(r, 18, ACCENT, FONT, bold=True)
        r2 = p.add_run(); r2.text = pt; _set_font(r2, 18, INK, FONT)
    footer(s)
    notes(s, "Four crisp takeaways \u2014 deliver with energy; this is what they remember. (~30s)")


# ─── Thank you ────────────────────────────────────────────────────────────────
def slide_thanks():
    s = add_slide()
    rect(s, 0, 0, EMU_W, EMU_H, fill=ACCENTDK)
    rect(s, MARGIN, Inches(2.7), Inches(0.9), Inches(0.14), fill=ACCENT)
    tb, tf = textbox(s, MARGIN, Inches(2.95), Inches(11.5), Inches(2))
    para(tf, "Thank you", 46, WHITE, name=FONT_SB, bold=True, first=True, space_after=8)
    para(tf, "Questions & feedback welcome", 20, RGBColor(0xB6,0xC2,0xD9))
    tb2, tf2 = textbox(s, MARGIN, Inches(5.3), Inches(11.5), Inches(1))
    para(tf2, "Rafael Aza  ·  Master Colloquium B  ·  Prof. Dr.-Ing. Marcus Purat",
         15, RGBColor(0x8F,0xA8,0xE0), first=True)
    notes(s, "Invite feedback explicitly \u2014 it shapes the final experiment and report. "
              "Keep Q&A answers short. (~30s)")


# ─── Build ────────────────────────────────────────────────────────────────────
slide_title()
slide_agenda()
slide_context_what()
slide_context_why()
slide_problem()
slide_requirements()
slide_architecture()
slide_design_highlights()
slide_terraform()
slide_methodology()
slide_results()
slide_status()
slide_outlook()
slide_summary()
slide_thanks()

out = os.path.join(os.path.dirname(__file__), "..", "docs", "presentation", "presentation.pptx")
out = os.path.abspath(out)
prs.save(out)
print(f"Saved {len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")
